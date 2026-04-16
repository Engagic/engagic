"""Document extractor using PyMuPDF with OCR fallback

Moved from: infocore/processing/pdf_extractor.py

Supports: PDF, DOCX (via PyMuPDF), legacy .doc (via antiword), RTF (via striprtf),
PPTX (via python-pptx).

Legislative formatting detection (strikethrough/underline):
- Detects thin filled rectangles (MS Word/LibreOffice export format)
- Strikethrough = deletions from law (line through middle of text)
- Underline = additions to law (line below text)
- Outputs as [DELETED: text] and [ADDED: text] tags in markdown
"""

import io
import os
import re
import subprocess
import tempfile
import time
import warnings
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Dict, Any, List, Tuple, Optional

import fitz  # PyMuPDF
import pytesseract
import requests
from PIL import Image

from config import get_logger
from exceptions import ExtractionError

logger = get_logger(__name__).bind(component="parser")

# Browser-like headers to avoid bot detection
DEFAULT_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36",
    "Accept": "application/pdf,application/octet-stream,*/*",
    "Accept-Language": "en-US,en;q=0.9",
    "Accept-Encoding": "gzip, deflate, br",
    "Connection": "keep-alive",
}

# Set conservative limit for OCR on scanned PDFs
# 100MP = ~300MB peak RAM (conservative for 2GB VPS with other services)
# Convert PIL warnings to errors to catch decompression bombs
Image.MAX_IMAGE_PIXELS = 100000000
warnings.simplefilter('error', Image.DecompressionBombWarning)

# Magic bytes for file format detection
_OLE2_MAGIC = b'\xd0\xcf\x11\xe0'    # Legacy .doc, .xls, .ppt (OLE2 Compound)
_ZIP_MAGIC = b'PK\x03\x04'           # .docx, .xlsx, .pptx (OOXML/ZIP)
_PDF_MAGIC = b'%PDF-'
_RTF_MAGIC = b'{\\rtf'


def _detect_format(data: bytes) -> str:
    """Detect document format from magic bytes.

    Returns 'pdf', 'docx', 'pptx', 'doc', 'rtf', or 'unknown'.
    For ZIP-based OOXML, peeks at [Content_Types].xml to distinguish
    PPTX (presentationml) from DOCX (wordprocessingml).
    """
    if data[:5] == _PDF_MAGIC:
        return "pdf"
    if data[:4] == _ZIP_MAGIC:
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                ct = zf.read("[Content_Types].xml").decode("utf-8", errors="ignore")
                if "presentationml" in ct:
                    return "pptx"
        except Exception:
            pass
        return "docx"
    if data[:4] == _OLE2_MAGIC:
        return "doc"
    if data[:5] == _RTF_MAGIC:
        return "rtf"
    return "unknown"


def _extract_legacy_doc(data: bytes) -> Optional[str]:
    """Extract text from legacy .doc (OLE2) using antiword. Returns text or None."""
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp_path = tmp.name
            tmp.write(data)
        result = subprocess.run(
            ["antiword", tmp_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
        logger.debug("antiword returned no text", exit_code=result.returncode, stderr=result.stderr[:200])
        return None
    except FileNotFoundError:
        logger.warning("antiword not installed, cannot extract legacy .doc")
        return None
    except subprocess.TimeoutExpired:
        logger.warning("antiword timed out")
        return None
    except Exception as e:
        logger.debug("antiword extraction failed", error=str(e))
        return None
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def _extract_rtf(data: bytes) -> Optional[str]:
    """Extract text from RTF using striprtf. Returns text or None."""
    try:
        from striprtf.striprtf import rtf_to_text
        rtf_str = data.decode("utf-8", errors="replace")
        text = rtf_to_text(rtf_str)
        return text if text and text.strip() else None
    except Exception as e:
        logger.debug("rtf extraction failed", error=str(e))
        return None


def _extract_pptx(data: bytes) -> Optional[str]:
    """Extract text from PPTX using python-pptx. Returns text or None."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
        return "\n".join(parts) if parts else None
    except Exception as e:
        logger.debug("pptx extraction failed", error=str(e))
        return None


def _detect_horizontal_lines(page: fitz.Page) -> List[Tuple[float, float, float]]:
    """
    Detect horizontal lines from drawing instructions.

    Strikethrough/underline in MS Word/LibreOffice are rendered as THIN FILLED RECTANGLES.

    Returns:
        list of (x0, x1, y_position) tuples for horizontal bars
    """
    lines = []
    paths = page.get_drawings()

    for path in paths:
        # Check if this path is a filled black rectangle (potential line)
        fill_color = path.get("fill")
        if fill_color and fill_color == (0.0, 0.0, 0.0):  # Black fill
            for item in path["items"]:
                if item[0] == "re":  # Rectangle
                    rect = item[1]
                    x0, y0, x1, y1 = rect

                    height = abs(y1 - y0)
                    width = abs(x1 - x0)

                    # Thin horizontal rectangle (height < 2 points, width > 5 points)
                    if height < 2 and width > 5:
                        y_mid = (y0 + y1) / 2
                        lines.append((x0, x1, y_mid))

                # Also check for actual line items (just in case)
                elif item[0] == "l":
                    p1, p2 = item[1:]
                    if abs(p1.y - p2.y) < 0.5:
                        line_x0 = min(p1.x, p2.x)
                        line_x1 = max(p1.x, p2.x)
                        lines.append((line_x0, line_x1, p1.y))

    return lines


def _match_lines_to_text(page: fitz.Page, lines: List[Tuple[float, float, float]]) -> Dict[str, List[Tuple[str, Tuple]]]:
    """
    Match horizontal lines to text spans and classify as strikethrough/underline.

    Returns dict with:
        'strikethrough': list of (text, bbox) tuples
        'underline': list of (text, bbox) tuples
    """
    result = {
        'strikethrough': [],
        'underline': []
    }

    # Get text with detailed positioning
    text_dict = page.get_text("dict")  # type: ignore[attr-defined]

    for block in text_dict["blocks"]:
        if block.get("type") != 0:  # Skip non-text blocks
            continue

        for line_obj in block.get("lines", []):
            for span in line_obj.get("spans", []):
                text = span.get("text", "").strip()
                if not text:
                    continue

                bbox = span["bbox"]  # (x0, y0, x1, y1)
                text_x0, text_y0, text_x1, text_y1 = bbox

                # Text vertical center and bottom
                text_bottom_y = text_y1
                text_height = text_y1 - text_y0

                # Check each horizontal line
                for line_x0, line_x1, line_y in lines:
                    # Check if line horizontally overlaps with text
                    # (with small tolerance for slight misalignment)
                    if not (line_x1 < text_x0 - 2 or line_x0 > text_x1 + 2):
                        # Line overlaps text horizontally

                        # Strikethrough: line passes through middle of text
                        # (within 30% to 70% of text height from top)
                        if text_y0 + 0.3 * text_height <= line_y <= text_y0 + 0.7 * text_height:
                            result['strikethrough'].append((text, bbox))
                            break

                        # Underline: line is just below text
                        # (within 3 points below text bottom, or just touching)
                        elif text_bottom_y - 1 <= line_y <= text_bottom_y + 3:
                            result['underline'].append((text, bbox))
                            break

    return result


def _has_legislative_legend(doc: fitz.Document, proximity_chars: int = 200, max_pages: int = 5) -> bool:
    """Check if document contains legislative formatting legend (clustered keywords).

    A true legislative legend has all 4 keyword types appearing close together,
    like: "Additions shown as underline, deletions shown as strikethrough"

    False positives occur when keywords are scattered throughout ordinance text
    describing what amendments do (e.g., "the addition of Section 12-345...").

    Legislative legends typically appear on the first few pages (cover or TOC),
    so we limit search to max_pages for performance.

    Args:
        doc: PyMuPDF document
        proximity_chars: Maximum distance between keywords to consider clustered (default 200)
        max_pages: Maximum pages to search (default 5 - legends are in first few pages)

    Returns:
        True only if all 4 keyword types appear within proximity_chars of each other
    """
    # Keyword patterns for each category
    addition_pattern = re.compile(r'\b(addition|added)\b')
    deletion_pattern = re.compile(r'\b(deletion|deleted)\b')
    underline_pattern = re.compile(r'\bunderline\b')
    strikethrough_pattern = re.compile(r'\bstrikethrough\b')

    # Search only first max_pages (legends appear early in documents)
    pages_to_search = min(len(doc), max_pages)
    for page_num in range(pages_to_search):
        text = doc[page_num].get_text().lower()  # type: ignore[attr-defined]

        # Find all positions of each keyword type
        addition_positions = [m.start() for m in addition_pattern.finditer(text)]
        deletion_positions = [m.start() for m in deletion_pattern.finditer(text)]
        underline_positions = [m.start() for m in underline_pattern.finditer(text)]
        strikethrough_positions = [m.start() for m in strikethrough_pattern.finditer(text)]

        # All 4 types must be present on this page
        if not (addition_positions and deletion_positions and
                underline_positions and strikethrough_positions):
            continue

        # Check if any combination of positions clusters within proximity_chars
        # Use underline as anchor since it's least likely to appear in normal text
        for u_pos in underline_positions:
            window_start = u_pos - proximity_chars
            window_end = u_pos + proximity_chars

            has_addition = any(window_start <= p <= window_end for p in addition_positions)
            has_deletion = any(window_start <= p <= window_end for p in deletion_positions)
            has_strikethrough = any(window_start <= p <= window_end for p in strikethrough_positions)

            if has_addition and has_deletion and has_strikethrough:
                return True

    return False


def _extract_text_with_formatting(page: fitz.Page, page_num: int) -> str:
    """
    Extract text from page with legislative formatting tags.

    Returns text with [DELETED: ...] and [ADDED: ...] tags for strikethrough/underline.
    """
    # Detect horizontal lines (strikethrough/underline indicators)
    lines = _detect_horizontal_lines(page)

    if not lines:
        # No formatting detected, return plain text
        return page.get_text(sort=True)  # type: ignore[attr-defined]

    # Match lines to text
    matched = _match_lines_to_text(page, lines)
    strikethrough_spans = {bbox: text for text, bbox in matched['strikethrough']}
    underline_spans = {bbox: text for text, bbox in matched['underline']}

    # Get text with detailed positioning
    text_dict = page.get_text("dict")  # type: ignore[attr-defined]
    result_parts = []

    for block in text_dict["blocks"]:
        if block.get("type") != 0:  # Skip non-text blocks
            continue

        block_parts = []
        for line_obj in block.get("lines", []):
            line_parts = []
            for span in line_obj.get("spans", []):
                text = span.get("text", "")
                bbox = tuple(span["bbox"])

                # Check if this span has formatting
                if bbox in strikethrough_spans:
                    line_parts.append(f"[DELETED: {text}]")
                elif bbox in underline_spans:
                    line_parts.append(f"[ADDED: {text}]")
                else:
                    line_parts.append(text)

            block_parts.append("".join(line_parts))

        result_parts.append("\n".join(block_parts))

    return "\n\n".join(result_parts)


class PdfExtractor:
    """PDF extractor using PyMuPDF with OCR fallback"""

    _instance: "PdfExtractor | None" = None

    @classmethod
    def shared(cls) -> "PdfExtractor":
        """Return a shared singleton instance."""
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    def __init__(self, ocr_threshold: int = 100, ocr_dpi: int = 200, detect_legislative_formatting: bool = True, max_ocr_workers: int | None = None):
        """Initialize PDF extractor

        Args:
            ocr_threshold: Minimum characters per page before triggering OCR fallback
            ocr_dpi: DPI for image rendering when using OCR (higher = better quality, slower)
                    Default 200 - balances quality vs memory on multi-core VPS
            detect_legislative_formatting: If True, detect strikethrough (deletions) and underline (additions)
                    in legislative documents with formatting legends and tag them as [DELETED: ...] and [ADDED: ...]
                    Only activates if document contains legislative formatting legend. Default: True (safe for all PDFs)
            max_ocr_workers: Maximum parallel OCR workers. Default: CPU count (min 1, max 4).
                    OCR is CPU-bound so more workers than cores causes thrashing.
        """
        self.ocr_threshold = ocr_threshold
        self.ocr_dpi = ocr_dpi
        self.detect_legislative_formatting = detect_legislative_formatting

        # Auto-detect optimal worker count based on CPU cores
        # OCR is CPU-bound, so workers > cores = context switch overhead
        if max_ocr_workers is None:
            cpu_count = os.cpu_count() or 1
            self.max_ocr_workers = min(cpu_count, 4)  # Cap at 4 for memory
        else:
            self.max_ocr_workers = max_ocr_workers

        # Prevent Tesseract internal threading when running multiple workers
        # Each worker gets 1 thread to avoid CPU thrashing
        os.environ.setdefault('OMP_THREAD_LIMIT', '1')

        logger.info(
            "PDF extractor initialized",
            ocr_workers=self.max_ocr_workers,
            ocr_dpi=self.ocr_dpi,
            cpu_count=os.cpu_count()
        )

    def _render_page_for_ocr(self, page) -> Optional[Tuple[bytes, int, int]]:
        """Render page to PNG bytes for OCR (main thread, not thread-safe)

        Args:
            page: PyMuPDF page object

        Returns:
            Tuple of (png_bytes, width, height) or None if too large
        """
        try:
            pix = page.get_pixmap(dpi=self.ocr_dpi)
            megapixels = (pix.width * pix.height) / 1000000

            logger.debug(
                "rendering page for OCR",
                page_num=page.number + 1,
                dpi=self.ocr_dpi,
                width=pix.width,
                height=pix.height,
                megapixels=round(megapixels, 1)
            )

            if megapixels > 100:
                logger.warning(
                    "page image too large for OCR",
                    page_num=page.number + 1,
                    megapixels=round(megapixels, 1)
                )
                return None

            return (pix.tobytes("png"), pix.width, pix.height)
        except Exception as e:  # Intentionally broad: graceful degradation, returns None
            logger.error("failed to render page", page_num=page.number + 1, error=str(e))
            return None

    def _ocr_from_bytes(self, png_bytes: bytes, page_num: int) -> str:
        """Run OCR on pre-rendered PNG bytes (thread-safe)

        Args:
            png_bytes: PNG image bytes
            page_num: Page number (1-indexed, for logging)

        Returns:
            Extracted text from OCR, or empty string on failure
        """
        try:
            img = Image.open(io.BytesIO(png_bytes)).convert('L')  # Grayscale
            # --oem 1: LSTM-only (2-3x faster than legacy+LSTM default)
            # --psm 3: Fully automatic page segmentation
            # timeout: Hard 60s cap per page to prevent hangs
            text = pytesseract.image_to_string(
                img,
                config='--oem 1 --psm 3',
                timeout=60
            )
            return text
        except (Image.DecompressionBombError, Image.DecompressionBombWarning):
            logger.warning("page image too large for OCR", page_num=page_num)
            return ""
        except RuntimeError as e:
            # pytesseract raises RuntimeError on timeout
            if "Tesseract process timeout" in str(e):
                logger.warning("OCR timeout on page", page_num=page_num)
            else:
                logger.error("OCR runtime error", page_num=page_num, error=str(e))
            return ""
        except (OSError, pytesseract.TesseractError) as e:
            logger.error("OCR failed", page_num=page_num, error=str(e), error_type=type(e).__name__)
            return ""

    def _ocr_page(self, page) -> str:
        """Extract text from page using OCR (sequential fallback)

        Args:
            page: PyMuPDF page object

        Returns:
            Extracted text from OCR, or empty string if image too large
        """
        rendered = self._render_page_for_ocr(page)
        if rendered is None:
            return ""
        png_bytes, _, _ = rendered
        return self._ocr_from_bytes(png_bytes, page.number + 1)

    def _ocr_pages_parallel(self, ocr_tasks: List[Tuple[int, bytes, str]]) -> Dict[int, str]:
        """Run OCR on multiple pages in parallel

        Args:
            ocr_tasks: List of (page_num, png_bytes, original_text) tuples

        Returns:
            Dict mapping page_num to OCR result text
        """
        if not ocr_tasks:
            return {}

        results = {}
        workers = min(self.max_ocr_workers, len(ocr_tasks))

        logger.info(
            "starting parallel OCR",
            pages=len(ocr_tasks),
            workers=workers
        )

        ocr_start = time.time()

        # Total OCR budget: 5 minutes for all pages combined
        # Prevents infinite hangs if Tesseract locks up on any page
        ocr_total_timeout = 300

        with ThreadPoolExecutor(max_workers=workers) as executor:
            # Submit all OCR jobs
            future_to_page = {
                executor.submit(self._ocr_from_bytes, png_bytes, page_num): (page_num, original_text)
                for page_num, png_bytes, original_text in ocr_tasks
            }

            # Collect results as they complete (with total timeout)
            completed_futures = set()
            try:
                for future in as_completed(future_to_page, timeout=ocr_total_timeout):
                    completed_futures.add(future)
                    page_num, original_text = future_to_page[future]
                    try:
                        # Per-page timeout as secondary safeguard
                        ocr_result = future.result(timeout=120)

                        # Decide whether to use OCR or keep original
                        if self._is_ocr_better(original_text, ocr_result, page_num):
                            results[page_num] = ocr_result
                        else:
                            results[page_num] = original_text

                    except TimeoutError:
                        logger.warning("OCR timeout on page", page_num=page_num)
                        results[page_num] = original_text
                    except Exception as e:  # Intentionally broad: catch any thread exception
                        logger.error("parallel OCR failed", page_num=page_num, error=str(e) or type(e).__name__)
                        results[page_num] = original_text

            except TimeoutError:
                # Total OCR budget exceeded - cancel remaining futures and use original text
                timed_out_pages = []
                for future, (page_num, original_text) in future_to_page.items():
                    if future not in completed_futures:
                        future.cancel()
                        results[page_num] = original_text
                        timed_out_pages.append(page_num)
                logger.warning(
                    "OCR total timeout exceeded, skipping remaining pages",
                    timed_out_pages=timed_out_pages,
                    timeout_seconds=ocr_total_timeout
                )

        ocr_time = time.time() - ocr_start
        logger.info(
            "parallel OCR complete",
            pages=len(ocr_tasks),
            ocr_time=round(ocr_time, 2),
            avg_per_page=round(ocr_time / len(ocr_tasks), 2) if ocr_tasks else 0
        )

        return results

    def _is_ocr_better(self, original: str, ocr_result: str, page_num: int) -> bool:
        """Determine if OCR result is better than original text

        Args:
            original: Original extracted text
            ocr_result: OCR-produced text
            page_num: Page number (1-indexed, for logging)

        Returns:
            True if OCR should be used, False if original should be kept
        """
        orig_chars = len(original.strip())
        ocr_chars = len(ocr_result.strip())

        # If OCR produced nothing, always keep original
        if ocr_chars == 0:
            logger.info("keeping original text, OCR produced nothing", page_num=page_num)
            return False

        # Calculate quality metrics for OCR
        ocr_letters = sum(1 for c in ocr_result if c.isalpha())
        ocr_letter_ratio = ocr_letters / len(ocr_result) if len(ocr_result) > 0 else 0
        ocr_words = len(ocr_result.split())

        # OCR is better if:
        # 1. Produced significantly more text (2x+ characters) with reasonable quality (>40% letters)
        # 2. OR produced more text with high quality (>70% letters)
        significantly_more = ocr_chars >= (orig_chars * 2) and ocr_letter_ratio > 0.4
        high_quality_improvement = ocr_chars > orig_chars and ocr_letter_ratio > 0.7

        if significantly_more or high_quality_improvement:
            logger.info(
                f"[PyMuPDF] Page {page_num}: Using OCR "
                f"({ocr_chars} chars, {ocr_words} words, {ocr_letter_ratio:.1%} letters > original {orig_chars} chars)"
            )
            return True
        else:
            logger.info(
                f"[PyMuPDF] Page {page_num}: Keeping original "
                f"(OCR: {ocr_chars} chars, {ocr_letter_ratio:.1%} letters not better than original {orig_chars} chars)"
            )
            return False

    def _extract_from_document(self, doc: fitz.Document, extract_links: bool, start_time: float) -> Dict[str, Any]:
        """Core extraction logic for opened PDF document

        Args:
            doc: Opened PyMuPDF document
            extract_links: Whether to extract hyperlinks
            start_time: Extraction start time (for timing)

        Returns:
            Dict with extraction results
        """
        page_texts = {}  # page_num -> text
        all_links = []
        ocr_tasks = []  # List of (page_num, png_bytes, original_text)

        # Check for legislative legend once (if formatting detection enabled)
        use_formatting = self.detect_legislative_formatting and _has_legislative_legend(doc)
        if use_formatting:
            logger.info("[PyMuPDF] Legislative formatting detected - tagging additions/deletions")

        # Per-page sanity cap. A single PDF page should never yield more than
        # a few KB of real text. When PyMuPDF sees a broken font CMap or a
        # mis-tagged content stream, get_text() can return megabytes of
        # gibberish for one page and the total exceeds the 1GB RLIMIT_AS
        # child budget or blows up multiprocessing pickle. Truncate early so
        # the rest of the document still extracts.
        _MAX_PAGE_CHARS = 200_000

        # Pass 1: Extract text from all pages, collect OCR tasks
        for page_num in range(len(doc)):
            page = doc[page_num]

            # Extract text (with or without legislative formatting detection)
            if use_formatting:
                page_text = _extract_text_with_formatting(page, page_num + 1)
            else:
                page_text = page.get_text(sort=True)  # type: ignore[attr-defined]

            if len(page_text) > _MAX_PAGE_CHARS:
                logger.warning(
                    "[PyMuPDF] Page yielded suspicious text volume, truncating",
                    page_num=page_num + 1,
                    chars=len(page_text),
                    limit=_MAX_PAGE_CHARS,
                )
                page_text = page_text[:_MAX_PAGE_CHARS]

            initial_char_count = len(page_text.strip())

            # If page has minimal text, queue for OCR
            if initial_char_count < self.ocr_threshold:
                logger.debug(
                    "page queued for OCR",
                    page_num=page_num + 1,
                    char_count=initial_char_count,
                    threshold=self.ocr_threshold
                )
                # Render page to PNG (main thread, PyMuPDF not thread-safe)
                rendered = self._render_page_for_ocr(page)
                if rendered:
                    png_bytes, _, _ = rendered
                    ocr_tasks.append((page_num + 1, png_bytes, page_text))
                else:
                    # Rendering failed, keep original
                    page_texts[page_num + 1] = page_text
            else:
                page_texts[page_num + 1] = page_text

            # Extract links if requested
            if extract_links:
                page_links = page.get_links()  # type: ignore[attr-defined]
                for link in page_links:
                    if 'uri' in link and link['uri']:
                        all_links.append({
                            'page': page_num + 1,
                            'url': link['uri'],
                            'rect': link.get('from', None),
                        })

        page_count = len(doc)

        # Pass 2: Run OCR in parallel (outside doc context, PNG bytes already captured)
        if ocr_tasks:
            ocr_results = self._ocr_pages_parallel(ocr_tasks)
            page_texts.update(ocr_results)

        # Count OCR pages (pages where OCR was actually used, not just attempted)
        ocr_pages = sum(
            1 for page_num, _, original in ocr_tasks
            if page_num in page_texts and page_texts[page_num] != original
        )

        # Assemble final text in page order
        text_parts = [
            f"--- PAGE {page_num} ---\n{page_texts[page_num]}"
            for page_num in sorted(page_texts.keys())
        ]
        full_text = "\n\n".join(text_parts)

        extraction_time = time.time() - start_time

        # Determine extraction method
        method = "pymupdf+ocr" if ocr_pages > 0 else "pymupdf"

        log_msg = f"[PyMuPDF] Extracted {page_count} pages, {len(full_text)} chars"
        if ocr_pages > 0:
            log_msg += f" (OCR: {ocr_pages} pages)"
        if extract_links:
            log_msg += f", {len(all_links)} links"
        log_msg += f" in {extraction_time:.2f}s"
        logger.info(log_msg)

        result = {
            "success": True,
            "text": full_text,
            "method": method,
            "page_count": page_count,
            "extraction_time": extraction_time,
            "ocr_pages": ocr_pages,
        }

        if extract_links:
            result["links"] = all_links

        return result

    def extract_from_url(self, url: str, extract_links: bool = False) -> Dict[str, Any]:
        """Extract text and optionally links from PDF URL

        Args:
            url: PDF URL to extract from
            extract_links: Whether to extract hyperlinks (default False for backward compatibility)

        Returns dict with extraction results:
        {
            'success': bool,
            'text': str,
            'method': str,
            'page_count': int,
            'extraction_time': float,
            'links': list (if extract_links=True),
            'error': str (if failed)
        }
        """
        start_time = time.time()

        try:
            response = requests.get(url, headers=DEFAULT_HEADERS, timeout=30)
            response.raise_for_status()
            pdf_bytes = response.content

            with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                return self._extract_from_document(doc, extract_links, start_time)

        except Exception as e:  # Intentionally broad: API boundary, convert to typed error
            extraction_time = time.time() - start_time
            logger.error("[PyMuPDF] extraction failed", url=url[:100], error=str(e), error_type=type(e).__name__, extraction_time=round(extraction_time, 2))
            raise ExtractionError(
                f"PDF extraction failed after {extraction_time:.1f}s",
                document_url=url,
                document_type="pdf",
                original_error=e
            ) from e

    def extract_from_bytes(self, pdf_bytes: bytes, extract_links: bool = False) -> Dict[str, Any]:
        """Extract text from document bytes (PDF, DOCX, legacy .doc, RTF).

        Detects format from magic bytes and routes to the appropriate extractor.
        PDF and DOCX go through PyMuPDF; legacy .doc uses antiword; RTF uses striprtf.

        Args:
            pdf_bytes: Document content as bytes
            extract_links: If True, also extract hyperlinks (PDF/DOCX only)

        Returns dict with extraction results (same format as extract_from_url)
        """
        start_time = time.time()
        fmt = _detect_format(pdf_bytes)

        # Legacy .doc (OLE2) -- fitz can't handle this format
        if fmt == "doc":
            text = _extract_legacy_doc(pdf_bytes)
            if text:
                extraction_time = time.time() - start_time
                logger.info("extracted legacy .doc via antiword", chars=len(text), extraction_time=round(extraction_time, 2))
                return {
                    "success": True,
                    "text": text,
                    "method": "antiword",
                    "page_count": 0,
                    "extraction_time": extraction_time,
                }
            raise ExtractionError("Legacy .doc extraction failed (antiword unavailable or returned no text)", document_type="doc")

        # RTF
        if fmt == "rtf":
            text = _extract_rtf(pdf_bytes)
            if text:
                extraction_time = time.time() - start_time
                logger.info("extracted rtf", chars=len(text), extraction_time=round(extraction_time, 2))
                return {
                    "success": True,
                    "text": text,
                    "method": "striprtf",
                    "page_count": 0,
                    "extraction_time": extraction_time,
                }
            raise ExtractionError("RTF extraction failed", document_type="rtf")

        # PPTX
        if fmt == "pptx":
            text = _extract_pptx(pdf_bytes)
            if text:
                extraction_time = time.time() - start_time
                logger.info("extracted pptx via python-pptx", chars=len(text), extraction_time=round(extraction_time, 2))
                return {
                    "success": True,
                    "text": text,
                    "method": "python-pptx",
                    "page_count": 0,
                    "extraction_time": extraction_time,
                }
            raise ExtractionError("PPTX extraction failed", document_type="pptx")

        # PDF, DOCX, and unknown formats -- PyMuPDF handles all of these
        try:
            with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                return self._extract_from_document(doc, extract_links, start_time)

        except Exception as e:  # Intentionally broad: API boundary, convert to typed error
            extraction_time = time.time() - start_time
            logger.error("[PyMuPDF] extraction from bytes failed", format=fmt, error=str(e), error_type=type(e).__name__, extraction_time=round(extraction_time, 2))
            raise ExtractionError(
                f"Document extraction failed after {extraction_time:.1f}s",
                document_type=fmt,
                original_error=e
            ) from e

    def validate_text(self, text: str) -> bool:
        """Validate text quality - basic check for now"""
        # Simple validation: check if text is not empty and has reasonable content
        if not text or len(text) < 100:
            return False

        # Check if text has reasonable letter ratio
        letters = sum(1 for c in text if c.isalpha())
        if letters / len(text) < 0.3:
            return False

        return True
